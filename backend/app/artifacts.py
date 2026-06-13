from __future__ import annotations

import html
import re
import uuid
from dataclasses import dataclass, field
from datetime import datetime
from pathlib import Path
from typing import Any

from docx import Document
from docx.enum.text import WD_ALIGN_PARAGRAPH
from docx.shared import Inches, Pt, RGBColor
from fpdf import FPDF
from openpyxl import Workbook
from openpyxl.styles import Alignment, Font, PatternFill
from openpyxl.utils import get_column_letter
from pptx import Presentation
from pptx.dml.color import RGBColor as PptRGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches as PptInches
from pptx.util import Pt as PptPt
from openpyxl.drawing.image import Image as ExcelImage

from .config import BACKEND_DIR, Settings

LOGO_PATH = BACKEND_DIR / "static" / "logo.png"

ESPRESSO = "1D1712"
GOLD = "E8B44F"
TOMATO = "DB553D"
CREAM = "FFF5DF"


@dataclass
class ArtifactContent:
    title: str
    subtitle: str = ""
    sections: list[tuple[str, str]] = field(default_factory=list)
    columns: list[str] = field(default_factory=list)
    rows: list[list[Any]] = field(default_factory=list)
    sources: list[str] = field(default_factory=list)


def _filename(extension: str) -> str:
    stamp = datetime.now().strftime("%Y%m%d_%H%M%S")
    return f"artifact_{stamp}_{uuid.uuid4().hex[:6]}.{extension}"


def _files_dir() -> Path:
    path = BACKEND_DIR / "static" / "files"
    path.mkdir(parents=True, exist_ok=True)
    return path


def _ascii(value: Any) -> str:
    return str(value).encode("latin-1", "replace").decode("latin-1")


def artifact_url(settings: Settings, path: Path) -> str:
    if not path.exists():
        raise FileNotFoundError(path)
    return f"{settings.public_base_url}/files/{path.name}"


def render_inline_html(content: ArtifactContent) -> str:
    sections = "".join(
        (
            '<section class="slide">'
            f'<span class="eyebrow">Al Dente Company Brain</span>'
            f"<h2>{html.escape(title)}</h2>"
            f"<div class=\"copy\">{body}</div>"
            "</section>"
        )
        for title, body in content.sections
    )
    return f"""<!doctype html>
<html lang="en"><head><meta charset="utf-8"><meta name="viewport" content="width=device-width">
<style>
body{{margin:0;padding:28px;background:#130f0c;color:#fff5df;font:15px/1.5 Inter,system-ui,sans-serif}}
.deck{{display:grid;grid-template-columns:repeat(2,minmax(0,1fr));gap:18px;max-width:1100px;margin:auto}}
.hero{{grid-column:1/-1;padding:28px 32px;border-radius:24px;background:linear-gradient(135deg,#2b2119,#17110d);border:1px solid #5f4329}}
.hero h1{{font:700 clamp(30px,5vw,58px)/1.02 Georgia,serif;margin:8px 0;color:#f2c66d}}
.hero p{{color:#d7c9ae;margin:0;max-width:720px}}
.slide{{min-height:220px;padding:24px;border-radius:22px;background:#211914;border:1px solid #493426;box-shadow:0 18px 60px #0008}}
.slide h2{{font:700 26px/1.1 Georgia,serif;color:#f2c66d;margin:10px 0 18px}}
.eyebrow{{color:#db553d;text-transform:uppercase;letter-spacing:.16em;font-size:10px;font-weight:800}}
.copy{{color:#eee2cd}}.copy strong{{color:#fff}}.copy ul{{padding-left:19px}}.copy li{{margin:8px 0}}
.sources{{grid-column:1/-1;color:#978572;font-size:12px;padding:6px 4px}}
@media(max-width:760px){{.deck{{grid-template-columns:1fr}}.hero{{grid-column:auto}}}}
</style></head><body><main class="deck">
<header class="hero"><span class="eyebrow">Account intelligence</span>
<h1>{html.escape(content.title)}</h1><p>{html.escape(content.subtitle)}</p></header>
{sections}<footer class="sources">Sources: {html.escape(", ".join(content.sources))}</footer>
</main></body></html>"""


def render_inline_markdown(content: ArtifactContent) -> str:
    parts = [f"# {content.title}"]
    if content.subtitle:
        parts.append(content.subtitle)
    for title, body in content.sections:
        clean = re.sub(r"<li>(.*?)</li>", r"- \1", body, flags=re.I | re.S)
        clean = re.sub(r"</?(?:ul|ol|p|strong)>", "", clean, flags=re.I)
        clean = html.unescape(re.sub(r"<[^>]+>", "", clean)).strip()
        parts.extend((f"## {title}", clean))
    if content.columns and content.rows:
        parts.append("| " + " | ".join(map(str, content.columns)) + " |")
        parts.append("| " + " | ".join("---" for _ in content.columns) + " |")
        parts.extend("| " + " | ".join(map(str, row)) + " |" for row in content.rows)
    if content.sources:
        parts.append("Sources: " + ", ".join(content.sources))
    return "\n\n".join(part for part in parts if part)


def write_pdf(content: ArtifactContent) -> Path:
    path = _files_dir() / _filename("pdf")
    pdf = FPDF()
    pdf.set_auto_page_break(auto=True, margin=14)
    pdf.add_page()
    pdf.set_fill_color(29, 23, 18)
    pdf.rect(0, 0, 210, 42, "F")
    if LOGO_PATH.exists():
        from PIL import Image as PILImage
        img = PILImage.open(LOGO_PATH).convert("RGBA")
        bg = PILImage.new("RGBA", img.size, (29, 23, 18, 255))
        pdf_img = PILImage.alpha_composite(bg, img).convert("RGB")
        w_mm = 25.4 * (img.width / img.height)
        pdf.image(pdf_img, x=210 - w_mm - 5, y=5, h=25.4)
    pdf.set_text_color(232, 180, 79)
    pdf.set_font("Helvetica", "B", 22)
    pdf.set_xy(14, 13)
    pdf.cell(182, 10, _ascii(content.title))
    pdf.set_text_color(245, 237, 222)
    pdf.set_font("Helvetica", "", 9)
    pdf.set_xy(14, 28)
    pdf.multi_cell(182, 5, _ascii(content.subtitle))
    pdf.set_text_color(35, 29, 23)
    pdf.set_y(50)
    for title, body in content.sections:
        pdf.set_font("Helvetica", "B", 13)
        pdf.set_text_color(167, 86, 55)
        pdf.cell(0, 7, _ascii(title), new_x="LMARGIN", new_y="NEXT")
        pdf.set_font("Helvetica", "", 9)
        pdf.set_text_color(35, 29, 23)
        pdf.multi_cell(0, 5, _ascii(re.sub(r"<[^>]+>", " ", body)))
        pdf.ln(2)
    if content.columns and content.rows:
        usable = 182
        widths = [usable / len(content.columns)] * len(content.columns)
        pdf.set_font("Helvetica", "B", 8)
        pdf.set_fill_color(232, 180, 79)
        for index, column in enumerate(content.columns):
            pdf.cell(widths[index], 7, _ascii(column), border=1, fill=True)
        pdf.ln()
        pdf.set_font("Helvetica", "", 7)
        for row in content.rows:
            for index, value in enumerate(row):
                pdf.cell(widths[index], 6, _ascii(value)[:34], border=1)
            pdf.ln()
    pdf.set_y(-16)
    pdf.set_font("Helvetica", "I", 7)
    pdf.set_text_color(110, 100, 90)
    pdf.cell(0, 5, _ascii("Sources: " + ", ".join(content.sources)))
    pdf.output(path)
    return path


def write_xlsx(content: ArtifactContent) -> Path:
    path = _files_dir() / _filename("xlsx")
    workbook = Workbook()
    sheet = workbook.active
    sheet.title = "Report"
    sheet["A1"] = content.title
    sheet["A1"].font = Font(size=20, bold=True, color=GOLD)
    sheet["A1"].fill = PatternFill("solid", fgColor=ESPRESSO)
    sheet.merge_cells(start_row=1, start_column=1, end_row=1, end_column=max(len(content.columns), 4))
    sheet["A2"] = content.subtitle
    sheet.merge_cells(start_row=2, start_column=1, end_row=2, end_column=max(len(content.columns), 4))
    sheet["A2"].alignment = Alignment(wrap_text=True)
    header_row = 4
    if content.columns:
        for col, value in enumerate(content.columns, 1):
            cell = sheet.cell(header_row, col, value)
            cell.font = Font(bold=True, color=CREAM)
            cell.fill = PatternFill("solid", fgColor=TOMATO)
            cell.alignment = Alignment(horizontal="center")
        for row_index, row in enumerate(content.rows, header_row + 1):
            for col, value in enumerate(row, 1):
                sheet.cell(row_index, col, value)
        sheet.freeze_panes = f"A{header_row + 1}"
        sheet.auto_filter.ref = (
            f"A{header_row}:{get_column_letter(len(content.columns))}{header_row + len(content.rows)}"
        )
        for col in range(1, len(content.columns) + 1):
            values = [str(sheet.cell(row, col).value or "") for row in range(1, sheet.max_row + 1)]
            sheet.column_dimensions[get_column_letter(col)].width = min(max(map(len, values)) + 3, 42)
    notes = workbook.create_sheet("Notes")
    notes.append(["Section", "Detail"])
    for title, body in content.sections:
        notes.append([title, re.sub(r"<[^>]+>", " ", body)])
    notes.append(["Sources", ", ".join(content.sources)])
    notes.column_dimensions["A"].width = 28
    notes.column_dimensions["B"].width = 100
    notes.freeze_panes = "A2"
    if LOGO_PATH.exists():
        img = ExcelImage(LOGO_PATH)
        aspect = img.width / img.height
        img.height = 96
        img.width = int(96 * aspect)
        logo_col = get_column_letter(max(len(content.columns) + 2, 8))
        sheet.add_image(img, f"{logo_col}1")
    workbook.save(path)
    return path


def write_docx(content: ArtifactContent) -> Path:
    path = _files_dir() / _filename("docx")
    document = Document()
    section = document.sections[0]
    section.top_margin = Inches(0.65)
    section.bottom_margin = Inches(0.65)
    title = document.add_heading(content.title, 0)
    title.alignment = WD_ALIGN_PARAGRAPH.CENTER
    title.runs[0].font.color.rgb = RGBColor(167, 86, 55)
    subtitle = document.add_paragraph(content.subtitle)
    subtitle.alignment = WD_ALIGN_PARAGRAPH.CENTER
    for heading, body in content.sections:
        document.add_heading(heading, level=1)
        document.add_paragraph(re.sub(r"<[^>]+>", " ", body))
    if content.columns and content.rows:
        table = document.add_table(rows=1, cols=len(content.columns))
        table.style = "Light Shading Accent 1"
        for index, column in enumerate(content.columns):
            table.rows[0].cells[index].text = str(column)
        for row in content.rows:
            cells = table.add_row().cells
            for index, value in enumerate(row):
                cells[index].text = str(value)
    document.add_paragraph("Sources: " + ", ".join(content.sources)).runs[0].font.size = Pt(8)
    if LOGO_PATH.exists():
        header = section.header
        section.header_distance = Inches(0.2)
        header_para = header.paragraphs[0] if header.paragraphs else header.add_paragraph()
        header_para.alignment = WD_ALIGN_PARAGRAPH.RIGHT
        run = header_para.add_run()
        run.add_picture(str(LOGO_PATH), height=Inches(1.0))
    document.save(path)
    return path


def write_pptx(content: ArtifactContent) -> Path:
    path = _files_dir() / _filename("pptx")
    presentation = Presentation()
    presentation.slide_width = PptInches(13.333)
    presentation.slide_height = PptInches(7.5)
    slides = [("Al Dente Company Brain", content.title + "\n" + content.subtitle)] + content.sections[:4]
    for eyebrow, body in slides:
        slide = presentation.slides.add_slide(presentation.slide_layouts[6])
        background = slide.background.fill
        background.solid()
        background.fore_color.rgb = PptRGBColor(29, 23, 18)
        accent = slide.shapes.add_shape(1, PptInches(0), PptInches(0), PptInches(0.18), PptInches(7.5))
        accent.fill.solid()
        accent.fill.fore_color.rgb = PptRGBColor(219, 85, 61)
        accent.line.fill.background()
        box = slide.shapes.add_textbox(PptInches(0.8), PptInches(0.75), PptInches(11.7), PptInches(5.8))
        frame = box.text_frame
        frame.clear()
        p = frame.paragraphs[0]
        p.text = eyebrow.upper()
        p.font.size = PptPt(13)
        p.font.bold = True
        p.font.color.rgb = PptRGBColor(219, 85, 61)
        p.space_after = PptPt(18)
        p2 = frame.add_paragraph()
        p2.text = re.sub(r"<[^>]+>", " ", body)
        p2.font.size = PptPt(28 if len(slides) == 1 else 22)
        p2.font.color.rgb = PptRGBColor(255, 245, 223)
        p2.alignment = PP_ALIGN.LEFT
        source_box = slide.shapes.add_textbox(PptInches(0.8), PptInches(6.9), PptInches(11.7), PptInches(0.3))
        source_p = source_box.text_frame.paragraphs[0]
        source_p.text = "Sources: " + ", ".join(content.sources)
        source_p.font.size = PptPt(8)
        source_p.font.color.rgb = PptRGBColor(150, 132, 112)
        if LOGO_PATH.exists():
            from PIL import Image as PILImage
            img = PILImage.open(LOGO_PATH)
            aspect = img.width / img.height
            w_inches = 1.0 * aspect
            slide.shapes.add_picture(str(LOGO_PATH), left=PptInches(13.333 - w_inches - 0.2), top=PptInches(0.2), height=PptInches(1.0))
    presentation.save(path)
    return path


def write_binary(content: ArtifactContent, artifact_type: str) -> Path:
    writers = {
        "pdf": write_pdf,
        "xlsx": write_xlsx,
        "docx": write_docx,
        "pptx": write_pptx,
    }
    if artifact_type not in writers:
        raise ValueError(f"Unsupported binary artifact type: {artifact_type}")
    return writers[artifact_type](content)
